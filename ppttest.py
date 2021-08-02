from pptx import Presentation
from pptx.util import Inches
#from pptx.dml.color import RGBColor
#from pptx.enum.dml import MSO_THEME_COLOR
from pptx.util import Pt

prs = Presentation("show.pptx")  #选择pptx文件，路径按照需求做更改
blank_slide_layout = prs.slide_layouts[6]   #布局选用空页
slide = prs.slides.add_slide(blank_slide_layout)  #用空页布局创建一页幻灯片

left, top, width, height = Inches(1), Inches(1), Inches(5), Inches(5)
#left为距离ppt左侧13英寸，此时我的电脑上图片在ppt右侧能被看见一小部分
#top为距离ppt上侧1英寸
#width、height为图片的长宽、高，都为5英寸
#通过多次调整，大致感受图片的位置

#图片路径
img_path = r'1.JPG'
#在ppt指定位置按预设值添加图片
pic = slide.shapes.add_picture(img_path, left, top, width, height)

left, top, width, height = Inches(6),Inches(4),Inches(8),Inches(8)
tf = slide.shapes.add_textbox(left=left,top=top,width=width,height=height).text_frame
tf.paragraphs[0].text = "Hello"
p = tf.add_paragraph()
run = p.add_run()
run.text="World!"

f  = run.font
f.name="Arial"
f.size = Pt(38)

prs.save(r'show.pptx')
